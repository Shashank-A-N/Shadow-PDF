import os
import uuid
import fitz  # PyMuPDF
import zipfile
import io
from pptx import Presentation
from pptx.util import Inches
from flask import Flask, request, send_file, send_from_directory, jsonify

# Initialize the Flask application
app = Flask(__name__, static_folder='static')

# Define the folder to store uploaded files temporarily
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

@app.route('/')
def index():
    """Serves the index.html file from the 'static' directory."""
    return send_from_directory(app.static_folder, 'index.html')

@app.route('/convert', methods=['POST'])
def convert_multiple_files():
    """Handles multiple PDF uploads, converts them, and returns links."""
    files = request.files.getlist("pdfFiles")
    if not files or all(f.filename == '' for f in files):
        return jsonify({"error": "No files selected."}), 400

    batch_id = uuid.uuid4().hex
    batch_folder = os.path.join(app.config['UPLOAD_FOLDER'], batch_id)
    os.makedirs(batch_folder)
    
    converted_files = []

    for file in files:
        if file and file.filename.lower().endswith('.pdf'):
            original_filename = file.filename
            # Sanitize filename to prevent directory traversal issues
            safe_original_filename = "".join(c for c in original_filename if c.isalnum() or c in ('.', '_', '-')).rstrip()
            pdf_path = os.path.join(batch_folder, safe_original_filename)
            pptx_filename = os.path.splitext(safe_original_filename)[0] + '.pptx'
            pptx_path = os.path.join(batch_folder, pptx_filename)

            try:
                file.save(pdf_path)
                
                pdf_document = fitz.open(pdf_path)
                pres = Presentation()
                # Use a 16:9 aspect ratio, common for presentations
                pres.slide_width = Inches(16)
                pres.slide_height = Inches(9)

                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    # Use a higher DPI for better image quality
                    pix = page.get_pixmap(dpi=150)
                    image_bytes = pix.tobytes("png")
                    
                    img_path = os.path.join(batch_folder, f"temp_image_{page_num}.png")
                    with open(img_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    # Use a blank slide layout (layout index 6 is typically blank)
                    blank_slide_layout = pres.slide_layouts[6]
                    slide = pres.slides.add_slide(blank_slide_layout)
                    
                    # Add picture, ensuring it fits the slide dimensions
                    slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=pres.slide_width, height=pres.slide_height)
                    os.remove(img_path)

                pdf_document.close()
                pres.save(pptx_path)
                
                converted_files.append({
                    "original": original_filename,
                    "converted": pptx_filename
                })
                os.remove(pdf_path)

            except Exception as e:
                print(f"Error converting {original_filename}: {e}")
                continue
    
    if not converted_files:
        return jsonify({"error": "No valid PDF files were converted."}), 400
        
    return jsonify({
        "batch_id": batch_id,
        "files": converted_files
    })

@app.route('/download/<batch_id>/<filename>')
def download_file(batch_id, filename):
    """Serves an individual converted PPTX file for download."""
    directory = os.path.join(app.config['UPLOAD_FOLDER'], batch_id)
    return send_from_directory(directory, filename, as_attachment=True)

@app.route('/download-zip/<batch_id>')
def download_zip(batch_id):
    """Creates and serves a ZIP archive of all converted files in a batch."""
    batch_folder = os.path.join(app.config['UPLOAD_FOLDER'], batch_id)
    
    memory_file = io.BytesIO()
    with zipfile.ZipFile(memory_file, 'w', zipfile.ZIP_DEFLATED) as zf:
        for filename in os.listdir(batch_folder):
            if filename.lower().endswith('.pptx'):
                file_path = os.path.join(batch_folder, filename)
                zf.write(file_path, arcname=filename)

    memory_file.seek(0)
    
    return send_file(
        memory_file,
        download_name=f'converted_files_{batch_id}.zip',
        as_attachment=True,
        mimetype='application/zip'
    )

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000)))

