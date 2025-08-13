

# import os
# import sys
# import tempfile
# import shutil
# import subprocess
# from pathlib import Path
# from datetime import datetime
# import pythoncom
# import win32com.client

# from flask import Flask, render_template, request, send_file, jsonify, abort


import os
import sys
import tempfile
import shutil
import subprocess
from pathlib import Path
from datetime import datetime
from flask import Flask, render_template, request, send_file, jsonify, abort

# PDF -> PPT deps
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches


# # PDF -> PPT deps
# import fitz  # PyMuPDF
# from pptx import Presentation
# from pptx.util import Inches


app = Flask(__name__, static_folder="static", template_folder="templates")
app.config["MAX_CONTENT_LENGTH"] = 300 * 1024 * 1024  # 300 MB

PPT_EXTS = {".ppt", ".pptx"}
PDF_EXTS = {".pdf"}


# ---------------------------
# Engine detection (PPT->PDF)
# ---------------------------

def which_libreoffice():
    for cand in ("soffice", "libreoffice"):
        p = shutil.which(cand)
        if p:
            return p

    if os.name == "nt":
        for c in (
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ):
            if os.path.exists(c):
                return c

    if sys.platform == "darwin":
        c = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.exists(c):
            return c

    env = os.environ.get("LIBREOFFICE_PATH", "").strip()
    if env and os.path.exists(env):
        return env
    return None


def powerpoint_available():
    if os.name != "nt":
        return False
    try:
        import win32com.client  # noqa
        return True
    except Exception:
        return False


# ---------------------------
# PPT -> PDF engines
# ---------------------------

def convert_with_libreoffice(soffice_path, input_file, out_dir):
    input_file = str(Path(input_file).resolve())
    out_dir = str(Path(out_dir).resolve())
    os.makedirs(out_dir, exist_ok=True)

    cmd = [
        soffice_path,
        "--headless", "--nologo", "--nodefault", "--invisible", "--nofirststartwizard",
        "--convert-to", "pdf",
        input_file,
        "--outdir", out_dir,
    ]
    subprocess.run(cmd, check=True)

    stem = Path(input_file).with_suffix("").name
    expected = Path(out_dir) / f"{stem}.pdf"
    if expected.exists():
        return str(expected)
    candidates = sorted(Path(out_dir).glob(f"{stem}*.pdf"), key=lambda p: p.stat().st_mtime)
    if candidates:
        return str(candidates[-1])
    raise FileNotFoundError("LibreOffice reported success but PDF not found.")


def convert_with_powerpoint(input_file, out_dir):
    # Initialize COM for this thread, crucial for use in web servers like Flask
    pythoncom.CoInitialize()

    input_file = str(Path(input_file).resolve())
    out_dir = Path(out_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    stem = Path(input_file).with_suffix("").name
    out_pdf = out_dir / f"{stem}.pdf"

    pp = None
    pres = None
    try:
        pp = win32com.client.Dispatch("PowerPoint.Application")
        # Keep PowerPoint invisible during the conversion
        pres = pp.Presentations.Open(input_file, WithWindow=False)
        pres.SaveAs(str(out_pdf), 32)  # 32 = PDF format

    finally:
        # IMPORTANT: Always close and quit to avoid orphaned PowerPoint processes
        if pres is not None:
            pres.Close()
        if pp is not None:
            pp.Quit()
        # Uninitialize COM for this thread
        pythoncom.CoUninitialize()

    if not out_pdf.exists():
        raise FileNotFoundError(f"PowerPoint did not create the PDF file: {out_pdf}")
    return str(out_pdf)


def ppt_to_pdf(input_file, prefer="AUTO"):
    tmp_out = tempfile.mkdtemp(prefix="ppt2pdf_out_")
    lo = which_libreoffice()
    pp_ok = powerpoint_available()

    if prefer == "LIBREOFFICE":
        if not lo:
            raise RuntimeError("LibreOffice not found.")
        return convert_with_libreoffice(lo, input_file, tmp_out), "LibreOffice"

    if prefer == "POWERPOINT":
        if not pp_ok:
            raise RuntimeError("PowerPoint COM not available (install PowerPoint + pywin32).")
        return convert_with_powerpoint(input_file, tmp_out), "PowerPoint"

    # AUTO logic: Prefer PowerPoint on Windows if available, otherwise use LibreOffice
    if pp_ok:
        return convert_with_powerpoint(input_file, tmp_out), "PowerPoint"
    if lo:
        return convert_with_libreoffice(lo, input_file, tmp_out), "LibreOffice"

    raise RuntimeError("No conversion engine found. Install LibreOffice or (on Windows) PowerPoint+pywin32.")


# ---------------------------
# PDF -> PPT (image-based)
# ---------------------------

def pdf_to_ppt_image_based(input_pdf, scale=2.0):
    """
    Convert each PDF page to an image and place it on its own slide.
    - scale: 1.0 = 72 dpi, 2.0 = 144 dpi, 3.0 = 216 dpi (quality vs. file size)
    Returns path to generated .pptx.
    """
    input_pdf = str(Path(input_pdf).resolve())
    tmp_out = tempfile.mkdtemp(prefix="pdf2ppt_out_")
    out_pptx = str(Path(tmp_out) / (Path(input_pdf).stem + ".pptx"))

    doc = fitz.open(input_pdf)
    prs = Presentation()

    # For each page, we set the slide size to match the page size (in inches)
    for page in doc:
        rect = page.rect  # in points (1 pt = 1/72 inch)
        width_in = float(rect.width) / 72.0
        height_in = float(rect.height) / 72.0

        prs.slide_width = Inches(width_in)
        prs.slide_height = Inches(height_in)

        # Render page to image
        mat = fitz.Matrix(scale, scale)  # scale up for sharper image
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_path = str(Path(tmp_out) / f"page_{page.number+1}.png")
        pix.save(img_path)

        # Add slide with the image filling the entire slide
        blank = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    prs.save(out_pptx)
    doc.close()
    return out_pptx


# ---------------------------
# Routes
# ---------------------------

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/health")
def health():
    return jsonify({
        "status": "ok",
        "libreoffice": bool(which_libreoffice()),
        "powerpoint_com": bool(powerpoint_available()),
        "soffice_path": which_libreoffice(),
    })


@app.route("/convert", methods=["POST"])
def convert():
    """
    Handles both directions:
    - mode = PPT2PDF  (accepts .ppt/.pptx)   -> returns .pdf
    - mode = PDF2PPT  (accepts .pdf)         -> returns .pptx
    Optional for PPT2PDF: engine = AUTO | LIBREOFFICE | POWERPOINT
    """
    mode = (request.form.get("mode") or "PPT2PDF").upper()
    engine = (request.form.get("engine") or "AUTO").upper()

    if "file" not in request.files:
        abort(400, "No file field named 'file' in form data.")

    f = request.files["file"]
    if not f.filename:
        abort(400, "No file selected.")

    ext = Path(f.filename).suffix.lower()

    tmp_in = tempfile.mkdtemp(prefix="conv_in_")
    src_name = Path(f.filename).name
    src_path = str(Path(tmp_in) / src_name)
    f.save(src_path)

    try:
        if mode == "PPT2PDF":
            if ext not in PPT_EXTS:
                abort(400, "Please upload a .ppt or .pptx file for PPT2PDF.")
            pdf_path, engine_used = ppt_to_pdf(src_path, prefer=engine)
            out_name = f"{Path(src_name).stem} (via {engine_used}) {datetime.now().strftime('%Y-%m-%d_%H-%M')}.pdf"

            resp = send_file(pdf_path, as_attachment=True, download_name=out_name, mimetype="application/pdf")
            @resp.call_on_close
            def _cleanup():
                try:
                    shutil.rmtree(Path(pdf_path).parent, ignore_errors=True)
                    shutil.rmtree(tmp_in, ignore_errors=True)
                except Exception:
                    pass
            return resp

        elif mode == "PDF2PPT":
            if ext not in PDF_EXTS:
                abort(400, "Please upload a .pdf file for PDF2PPT.")
            pptx_path = pdf_to_ppt_image_based(src_path, scale=2.0)
            out_name = f"{Path(src_name).stem} (slides from PDF) {datetime.now().strftime('%Y-%m-%d_%H-%M')}.pptx"

            resp = send_file(pptx_path, as_attachment=True, download_name=out_name,
                             mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            @resp.call_on_close
            def _cleanup2():
                try:
                    shutil.rmtree(Path(pptx_path).parent, ignore_errors=True)
                    shutil.rmtree(tmp_in, ignore_errors=True)
                except Exception:
                    pass
            return resp

        else:
            abort(400, "Unknown mode. Use PPT2PDF or PDF2PPT.")

    except subprocess.CalledProcessError as e:
        shutil.rmtree(tmp_in, ignore_errors=True)
        abort(500, f"LibreOffice failed: {e}")
    except Exception as e:
        shutil.rmtree(tmp_in, ignore_errors=True)
        abort(500, f"Conversion failed: {e}")

def which_libreoffice():
    for cand in ("soffice", "libreoffice"):
        p = shutil.which(cand)
        if p:
            print(f"[DEBUG] Found LibreOffice binary: {p}")
            return p
    print("[DEBUG] LibreOffice not found in PATH.")
    ...



if __name__ == "__main__":
    app.run(host="127.0.0.1", port=5000, debug=True)


